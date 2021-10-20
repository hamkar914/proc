''''Test script written by Hampus Karlsson, August 2021.
    e-mail: hamka@chalmers.se/hkarlsson914@gmail.com
    put a bunch of deconvuluted/fitted 1D spectra
    as .png into a pptx.'''

import os
import pptx.util
from pptx.util import Inches

up_path = "C:/Users/hamka/Documents/data/"
data_fldr = "211019/5/"

ddir = up_path+data_fldr+"pdata/2/"
input_figs = [x for x in os.listdir(ddir) if x[-4:]==".png"]


prs = pptx.Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
slide_layout = prs.slide_layouts[6]

x_coords = [0.0,5.0,0.0,5.0,0.0,5.0,0.0,
            5.0,0.0,5.0,0.0,5.0,0.0,5.0,
            0.0,5.0,0.0,5.0]

y_coords = [0.25,0.25,0.25,0.25,0.25,0.25,
            0.25,0.25,0.25,0.25,0.25,0.25,
            0.25,0.25,0.25,0.25,0.25,0.25]

slide = prs.slides.add_slide(slide_layout)

for i in range(len(input_figs)):

    if i % 2 == 0:
        slide = prs.slides.add_slide(slide_layout)

    pic_path = up_path+data_fldr+"pdata/2/dcon_vd_"+str(i+1)+".png"
    picture = slide.shapes.add_picture(pic_path,Inches(x_coords[i]),Inches(y_coords[i]),width=Inches(5))

prs.save(up_path+data_fldr+"pdata/2/test.pptx")


